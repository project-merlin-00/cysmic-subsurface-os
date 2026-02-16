"""
Phase 4: Report Builder Service
PDF, DOCX, PPTX Report Generation
"""
import os
import json
import base64
from typing import List, Optional, Dict, Any
from datetime import datetime
from pathlib import Path
from io import BytesIO
from sqlalchemy.orm import Session

from backend.models.models import ReportTemplate, Report, Well, AnalysisResult
from backend.schemas import phase4 as schemas


class ReportBuilder:
    """Service for generating PDF, DOCX, and PPTX reports"""
    
    def __init__(self, db: Session):
        self.db = db
        self.base_path = Path("/home/arthur/.openclaw/workspace/cysmic-subsurface-os/data/reports")
        self.base_path.mkdir(parents=True, exist_ok=True)
    
    # ===================
    # Template Management
    # ===================
    
    def create_template(
        self,
        name: str,
        format: str,
        sections: List[dict],
        description: Optional[str] = None,
        owner_id: int = 1
    ) -> ReportTemplate:
        """Create a new report template"""
        template = ReportTemplate(
            name=name,
            description=description,
            format=format,
            sections=sections,
            owner_id=owner_id
        )
        self.db.add(template)
        self.db.commit()
        self.db.refresh(template)
        return template
    
    def get_templates(self, format: Optional[str] = None) -> List[ReportTemplate]:
        """Get all report templates"""
        query = self.db.query(ReportTemplate)
        if format:
            query = query.filter(ReportTemplate.format == format)
        return query.order_by(ReportTemplate.name).all()
    
    def get_template(self, template_id: int) -> Optional[ReportTemplate]:
        """Get a specific template"""
        return self.db.query(ReportTemplate).filter(
            ReportTemplate.id == template_id
        ).first()
    
    def get_default_template(self, format: str) -> Optional[ReportTemplate]:
        """Get default template for a format"""
        return self.db.query(ReportTemplate).filter(
            ReportTemplate.format == format,
            ReportTemplate.is_default == True
        ).first()
    
    # ===================
    # Report Generation
    # ===================
    
    def create_report(
        self,
        name: str,
        format: str,
        wells: List[int],
        analyses: List[int],
        sections: List[dict],
        title: str = "CYSMIC Report",
        subtitle: Optional[str] = None,
        template_id: Optional[int] = None,
        owner_id: int = 1
    ) -> Report:
        """Create and generate a report"""
        
        # Create report record
        report = Report(
            name=name,
            format=format,
            wells=wells,
            analyses=analyses,
            sections=sections,
            template_id=template_id,
            owner_id=owner_id,
            status="generating"
        )
        self.db.add(report)
        self.db.commit()
        self.db.refresh(report)
        
        # Generate the report
        try:
            file_path = self._generate_report(report, format, title, subtitle)
            report.file_path = file_path
            report.status = "completed"
            report.completed_at = datetime.utcnow()
        except Exception as e:
            report.status = "failed"
            print(f"Report generation failed: {e}")
        
        self.db.commit()
        self.db.refresh(report)
        
        return report
    
    def _generate_report(
        self,
        report: Report,
        format: str,
        title: str,
        subtitle: Optional[str]
    ) -> str:
        """Generate the actual report file"""
        
        if format == "pdf":
            return self._generate_pdf(report, title, subtitle)
        elif format == "docx":
            return self._generate_docx(report, title, subtitle)
        elif format == "pptx":
            return self._generate_pptx(report, title, subtitle)
        else:
            raise ValueError(f"Unsupported format: {format}")
    
    def _generate_pdf(self, report: Report, title: str, subtitle: Optional[str]) -> str:
        """Generate PDF report"""
        
        # Try to use reportlab if available, otherwise create a simple text-based PDF
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
            from reportlab.lib import colors
            
            # Create PDF
            pdf_path = self.base_path / f"report_{report.id}.pdf"
            doc = SimpleDocTemplate(str(pdf_path), pagesize=letter, topMargin=0.5*inch)
            
            styles = getSampleStyleSheet()
            story = []
            
            # Title
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                spaceAfter=12,
                textColor=colors.HexColor('#d44211')
            )
            story.append(Paragraph(title, title_style))
            
            if subtitle:
                story.append(Paragraph(subtitle, styles['Heading2']))
            
            story.append(Spacer(1, 0.3*inch))
            
            # Wells section
            wells = self.db.query(Well).filter(Well.id.in_(report.wells)).all()
            if wells:
                story.append(Paragraph("Wells Overview", styles['Heading2']))
                
                # Create wells table
                well_data = [['Well Name', 'Field', 'Status', 'Type', 'TVD (m)']]
                for well in wells:
                    well_data.append([
                        well.name or '-',
                        well.field or '-',
                        well.status.value if well.status else '-',
                        well.well_type or '-',
                        str(well.total_depth_tvd) if well.total_depth_tvd else '-'
                    ])
                
                t = Table(well_data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#d44211')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, -1), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.white),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(t)
                story.append(Spacer(1, 0.3*inch))
            
            # Analysis Results section
            analyses = self.db.query(AnalysisResult).filter(
                AnalysisResult.id.in_(report.analyses)
            ).all() if report.analyses else []
            
            if analyses:
                story.append(Paragraph("Analysis Results", styles['Heading2']))
                
                for analysis in analyses:
                    story.append(Paragraph(f"<b>{analysis.name or analysis.analysis_type}</b>", styles['Heading3']))
                    
                    # Add parameters
                    if analysis.parameters:
                        story.append(Paragraph("Parameters:", styles['Normal']))
                        for key, value in analysis.parameters.items():
                            story.append(Paragraph(f"  {key}: {value}", styles['Normal']))
                    
                    # Add results summary
                    if analysis.results:
                        story.append(Paragraph("Results:", styles['Normal']))
                        for key, value in analysis.results.items():
                            story.append(Paragraph(f"  {key}: {value}", styles['Normal']))
                    
                    story.append(Spacer(1, 0.2*inch))
            
            # Custom sections
            for section in report.sections:
                section_type = section.get('section_type')
                section_title = section.get('title', '')
                
                story.append(Paragraph(section_title, styles['Heading2']))
                
                if section_type == 'text':
                    content = section.get('content', '')
                    story.append(Paragraph(content, styles['Normal']))
                elif section_type == 'table':
                    table_data = section.get('content', {}).get('data', [])
                    if table_data:
                        t = Table(table_data)
                        t.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, -1), 9),
                            ('GRID', (0, 0), (-1, -1), 1, colors.black)
                        ]))
                        story.append(t)
                
                story.append(Spacer(1, 0.2*inch))
            
            # Build PDF
            doc.build(story)
            return str(pdf_path)
            
        except ImportError:
            # Fallback: create simple text file
            return self._generate_text_report(report, title, subtitle, "pdf")
    
    def _generate_docx(self, report: Report, title: str, subtitle: Optional[str]) -> str:
        """Generate DOCX report"""
        
        try:
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()
            
            # Title
            heading = doc.add_heading(title, 0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            if subtitle:
                p = doc.add_paragraph(subtitle)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_paragraph()
            
            # Wells section
            wells = self.db.query(Well).filter(Well.id.in_(report.wells)).all()
            if wells:
                doc.add_heading('Wells Overview', level=1)
                
                # Table
                table = doc.add_table(rows=1, cols=5)
                table.style = 'Table Grid'
                
                # Header row
                hdr_cells = table.rows[0].cells
                hdr_cells[0].text = 'Well Name'
                hdr_cells[1].text = 'Field'
                hdr_cells[2].text = 'Status'
                hdr_cells[3].text = 'Type'
                hdr_cells[4].text = 'TVD (m)'
                
                for well in wells:
                    row_cells = table.add_row().cells
                    row_cells[0].text = well.name or '-'
                    row_cells[1].text = well.field or '-'
                    row_cells[2].text = well.status.value if well.status else '-'
                    row_cells[3].text = well.well_type or '-'
                    row_cells[4].text = str(well.total_depth_tvd) if well.total_depth_tvd else '-'
                
                doc.add_paragraph()
            
            # Analysis Results
            analyses = self.db.query(AnalysisResult).filter(
                AnalysisResult.id.in_(report.analyses)
            ).all() if report.analyses else []
            
            if analyses:
                doc.add_heading('Analysis Results', level=1)
                
                for analysis in analyses:
                    doc.add_heading(analysis.name or analysis.analysis_type, level=2)
                    
                    if analysis.parameters:
                        doc.add_paragraph('Parameters:')
                        for key, value in analysis.parameters.items():
                            doc.add_paragraph(f'  {key}: {value}', style='List Bullet')
                    
                    if analysis.results:
                        doc.add_paragraph('Results:')
                        for key, value in analysis.results.items():
                            doc.add_paragraph(f'  {key}: {value}', style='List Bullet')
                    
                    doc.add_paragraph()
            
            # Custom sections
            for section in report.sections:
                section_type = section.get('section_type')
                section_title = section.get('title', '')
                
                doc.add_heading(section_title, level=1)
                
                if section_type == 'text':
                    content = section.get('content', '')
                    doc.add_paragraph(content)
                elif section_type == 'table':
                    table_data = section.get('content', {}).get('data', [])
                    if table_data and len(table_data) > 1:
                        table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                        table.style = 'Table Grid'
                        
                        for i, row in enumerate(table_data):
                            for j, cell in enumerate(row):
                                table.rows[i].cells[j].text = str(cell)
                
                doc.add_paragraph()
            
            # Save
            docx_path = self.base_path / f"report_{report.id}.docx"
            doc.save(str(docx_path))
            return str(docx_path)
            
        except ImportError:
            return self._generate_text_report(report, title, subtitle, "docx")
    
    def _generate_pptx(self, report: Report, title: str, subtitle: Optional[str]) -> str:
        """Generate PPTX presentation"""
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle_field = slide.placeholders[1]
            
            title.text = title
            if subtitle:
                subtitle_field.text = subtitle
            
            # Wells slide
            wells = self.db.query(Well).filter(Well.id.in_(report.wells)).all()
            if wells:
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
                
                # Add title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = "Wells Overview"
                p.font.size = Pt(36)
                p.font.bold = True
                
                # Add table
                rows = len(wells) + 1
                cols = 5
                table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(12), Inches(3)).table
                
                # Headers
                headers = ['Well Name', 'Field', 'Status', 'Type', 'TVD (m)']
                for i, h in enumerate(headers):
                    table.cell(0, i).text = h
                
                # Data
                for i, well in enumerate(wells):
                    table.cell(i+1, 0).text = well.name or '-'
                    table.cell(i+1, 1).text = well.field or '-'
                    table.cell(i+1, 2).text = well.status.value if well.status else '-'
                    table.cell(i+1, 3).text = well.well_type or '-'
                    table.cell(i+1, 4).text = str(well.total_depth_tvd) if well.total_depth_tvd else '-'
            
            # Analysis slides
            analyses = self.db.query(AnalysisResult).filter(
                AnalysisResult.id.in_(report.analyses)
            ).all() if report.analyses else []
            
            for analysis in analyses:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                
                # Title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = analysis.name or analysis.analysis_type
                p.font.size = Pt(32)
                p.font.bold = True
                
                # Content
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
                tf = content_box.text_frame
                tf.word_wrap = True
                
                if analysis.parameters:
                    p = tf.paragraphs[0]
                    p.text = "Parameters:"
                    p.font.size = Pt(18)
                    p.font.bold = True
                    
                    for key, value in analysis.parameters.items():
                        p = tf.add_paragraph()
                        p.text = f"  {key}: {value}"
                        p.font.size = Pt(16)
                
                if analysis.results:
                    p = tf.add_paragraph()
                    p.text = "Results:"
                    p.font.size = Pt(18)
                    p.font.bold = True
                    
                    for key, value in analysis.results.items():
                        p = tf.add_paragraph()
                        p.text = f"  {key}: {value}"
                        p.font.size = Pt(16)
            
            # Custom sections
            for section in report.sections:
                section_title = section.get('title', '')
                
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = section_title
                p.font.size = Pt(32)
                p.font.bold = True
            
            # Save
            pptx_path = self.base_path / f"report_{report.id}.pptx"
            prs.save(str(pptx_path))
            return str(pptx_path)
            
        except ImportError:
            return self._generate_text_report(report, title, subtitle, "pptx")
    
    def _generate_text_report(
        self,
        report: Report,
        title: str,
        subtitle: Optional[str],
        format: str
    ) -> str:
        """Fallback: Generate a simple text report"""
        
        lines = []
        lines.append("=" * 60)
        lines.append(title.upper())
        if subtitle:
            lines.append(subtitle)
        lines.append("=" * 60)
        lines.append(f"Generated: {datetime.utcnow().isoformat()}")
        lines.append("")
        
        # Wells
        wells = self.db.query(Well).filter(Well.id.in_(report.wells)).all()
        if wells:
            lines.append("WELLS OVERVIEW")
            lines.append("-" * 40)
            for well in wells:
                lines.append(f"  {well.name}")
                lines.append(f"    Field: {well.field or '-'}")
                lines.append(f"    Status: {well.status.value if well.status else '-'}")
                lines.append(f"    Type: {well.well_type or '-'}")
                lines.append(f"    TVD: {well.total_depth_tvd or '-'}")
                lines.append("")
        
        # Analyses
        analyses = self.db.query(AnalysisResult).filter(
            AnalysisResult.id.in_(report.analyses)
        ).all() if report.analyses else []
        
        if analyses:
            lines.append("ANALYSIS RESULTS")
            lines.append("-" * 40)
            for analysis in analyses:
                lines.append(f"  {analysis.name or analysis.analysis_type}")
                if analysis.parameters:
                    for key, value in analysis.parameters.items():
                        lines.append(f"    {key}: {value}")
                lines.append("")
        
        content = "\n".join(lines)
        
        # Save as text file (with requested extension for compatibility)
        ext = format if format in ['txt', 'md'] else 'txt'
        txt_path = self.base_path / f"report_{report.id}.{ext}"
        
        with open(txt_path, 'w') as f:
            f.write(content)
        
        return str(txt_path)
    
    def get_report(self, report_id: int) -> Optional[Report]:
        """Get a report by ID"""
        return self.db.query(Report).filter(Report.id == report_id).first()
    
    def get_user_reports(self, owner_id: int) -> List[Report]:
        """Get all reports for a user"""
        return self.db.query(Report).filter(
            Report.owner_id == owner_id
        ).order_by(Report.created_at.desc()).all()
    
    def delete_report(self, report_id: int, owner_id: int) -> bool:
        """Delete a report"""
        report = self.db.query(Report).filter(
            Report.id == report_id,
            Report.owner_id == owner_id
        ).first()
        
        if report:
            # Delete file if exists
            if report.file_path and os.path.exists(report.file_path):
                os.remove(report.file_path)
            
            self.db.delete(report)
            self.db.commit()
            return True
        
        return False
